from pptx import Presentation
from pptx.util import Inches, Pt

pptx_path = "毕业设计答辩-柜门板材切割排版系统.pptx"
prs = Presentation(pptx_path)

print(f"PPT文件信息：")
print(f"- 幻灯片数量: {len(prs.slides)}")
print(f"- 幻灯片尺寸: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
print()

for i, slide in enumerate(prs.slides, 1):
    print(f"=== 第{i}页 ===")

    # 提取文本内容
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            if len(text) > 50:
                text = text[:50] + "..."
            texts.append(text)

    if texts:
        for text in texts[:3]:  # 只显示前3个文本框
            print(f"  - {text}")
    else:
        print("  (无文本内容)")
    print()
